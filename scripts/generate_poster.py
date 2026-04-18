#!/usr/bin/env python3
"""
Generate scientific poster for X-Meeting 2026 in PPTX format.
Dimensions: 90 cm (width) x 110 cm (height)

v3: Fixed layout — larger fonts, constrained figure heights,
    better space distribution across columns.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

PROJECT_ROOT = Path(__file__).resolve().parent.parent
FIGURES_DIR = PROJECT_ROOT / "output" / "paper_figures"
OUTPUT = PROJECT_ROOT / "output" / "poster_xmeeting2026.pptx"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xA5)
LIGHT_BLUE = RGBColor(0xE3, 0xEF, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY = RGBColor(0x55, 0x55, 0x55)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_ALT = RGBColor(0xEC, 0xF0, 0xF6)

# Poster dimensions
W = Cm(90)
H = Cm(110)

# Layout
MARGIN = Cm(2)
COL_GAP = Cm(1.5)
CONTENT_W = Cm(86)
COL_W = (CONTENT_W - 2 * COL_GAP) / 3  # ~27.67 cm


def add_colored_box(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_section_header(slide, left, top, width, title):
    bar_h = Cm(1.3)
    add_colored_box(slide, left, top, width, bar_h, DARK_BLUE)
    txBox = slide.shapes.add_textbox(left + Cm(0.5), top + Cm(0.05),
                                     width - Cm(1), bar_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT
    return top + bar_h + Cm(0.3)


def add_subsection_header(slide, left, top, width, title):
    add_colored_box(slide, left, top, width, Cm(0.08), ACCENT_BLUE)
    txBox = slide.shapes.add_textbox(left, top + Cm(0.15), width, Cm(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT
    return top + Cm(1.4)


def add_body_text(slide, left, top, width, text, font_size=16,
                  color=DARK_GRAY):
    """Returns estimated height of text block."""
    # Estimate: chars per line based on font size and column width
    # col_w ~ 27.67 cm, at 16pt ~50 chars/line
    chars_per_line = max(30, int(width / Cm(1) * 1.8 * (16 / font_size)))
    total_lines = 0
    for paragraph in text.split('\n'):
        if paragraph.strip() == '':
            total_lines += 0.5
        else:
            total_lines += max(1, len(paragraph) / chars_per_line)
    est_height = Cm(total_lines * font_size * 0.048 + 0.5)

    txBox = slide.shapes.add_textbox(left, top, width, est_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, paragraph in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.JUSTIFY
    return est_height


def add_highlight_box(slide, left, top, width, text, font_size=15,
                      bg_color=LIGHT_BLUE, text_color=DARK_BLUE):
    """Highlighted box with background color."""
    chars_per_line = max(30, int(width / Cm(1) * 1.7 * (16 / font_size)))
    total_lines = 0
    for paragraph in text.split('\n'):
        if paragraph.strip() == '':
            total_lines += 0.3
        else:
            total_lines += max(1, len(paragraph) / chars_per_line)
    box_h = Cm(total_lines * font_size * 0.048 + 0.8)

    add_colored_box(slide, left, top, width, box_h, bg_color)
    txBox = slide.shapes.add_textbox(left + Cm(0.4), top + Cm(0.2),
                                     width - Cm(0.8), box_h - Cm(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, paragraph in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = paragraph
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Arial"
        p.font.bold = (i == 0)
        p.space_after = Pt(3)
        p.alignment = PP_ALIGN.LEFT
    return box_h


def add_figure(slide, left, top, img_path, width=None, max_height=None):
    """Add figure with optional max height constraint. Returns actual height."""
    path = Path(img_path)
    if not path.exists():
        print(f"  WARNING: not found: {img_path}")
        return Cm(3)

    # Get image dimensions to compute aspect ratio
    with Image.open(path) as img:
        img_w, img_h = img.size
    aspect = img_h / img_w

    if width is None:
        width = COL_W

    natural_h = int(width) * aspect  # in EMU
    if max_height is not None and natural_h > int(max_height):
        # Constrain by height, reduce width proportionally
        pic = slide.shapes.add_picture(str(path), left, top, height=max_height)
        # Center horizontally within column
        actual_w = pic.width
        offset = (int(width) - actual_w) // 2
        if offset > 0:
            pic.left = int(left) + offset
        return max_height
    else:
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        return pic.height


def add_caption(slide, left, top, width, text):
    txBox = slide.shapes.add_textbox(left, top, width, Cm(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = MED_GRAY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    return Cm(1.5)


def add_data_table(slide, left, top, width):
    rows = [
        ("Data Type", "Format", "Files", "Size"),
        ("ChIP-seq histones (6 marks \u00d7 11 lines)", "bigWig", "66", "~27.4 GB"),
        ("ChIP-seq CTCF (signal + peaks)", "bigWig/BED", "22", "~3.6 GB"),
        ("DNase-seq (signal + peaks)", "bigWig/BED", "22", "~5.5 GB"),
        ("DNA methylation (WGBS)", "BED", "8", "~4.6 GB"),
        ("Repli-seq (replication timing)", "bigWig", "6", "~0.05 GB"),
        ("Annotations (GENCODE, CpG, HKGs)", "GTF/BED/TSV", "6", "<0.1 GB"),
        ("Total", "", "130", "~41.2 GB"),
    ]
    n_rows = len(rows)
    n_cols = 4
    row_h = Cm(1.0)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top,
                                 width, row_h * n_rows).table
    # Column widths proportional
    tbl.columns[0].width = Cm(13.0)
    tbl.columns[1].width = Cm(5.0)
    tbl.columns[2].width = Cm(3.5)
    tbl.columns[3].width = Cm(5.0)

    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
            elif i == n_rows - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = DARK_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (TABLE_ROW_ALT if i % 2 == 0
                                            else WHITE)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = DARK_GRAY

    return row_h * n_rows + Cm(0.2)


def add_candidates_table(slide, left, top, width):
    rows = [
        ("Rank", "Genes", "Chr", "Stability"),
        ("#1", "SMIM27 / TOPORS", "9", "100%"),
        ("#2", "CNOT8 / FAXDC2", "5", "100%"),
        ("#3", "SEPTIN2 / HDLBP", "2", "100%"),
        ("#4", "RETREG2 / CNPPD1", "2", "100%"),
        ("#7", "ZBTB1 / ZBTB25", "14", "100%"),
        ("#8", "MICOS10 / CAPZB", "1", "100%"),
        ("#9", "AP4S1 / STRN3", "14", "100%"),
    ]
    n_rows = len(rows)
    n_cols = 4
    row_h = Cm(0.95)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top,
                                 width, row_h * n_rows).table
    tbl.columns[0].width = Cm(3.5)
    tbl.columns[1].width = Cm(12.5)
    tbl.columns[2].width = Cm(4.0)
    tbl.columns[3].width = Cm(5.0)

    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (TABLE_ROW_ALT if i % 2 == 0
                                            else WHITE)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = DARK_GRAY

    return row_h * n_rows + Cm(0.2)


# ================================================================
#  MAX FIGURE HEIGHTS — prevent any single figure from dominating
# ================================================================
MAX_FIG_H = Cm(18)        # default cap
MAX_FIG_H_SMALL = Cm(14)  # for smaller slots


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # ================================================================
    #  HEADER
    # ================================================================
    header_h = Cm(13)
    add_colored_box(slide, Cm(0), Cm(0), W, header_h, DARK_BLUE)
    add_colored_box(slide, Cm(0), header_h, W, Cm(0.3), GOLD)

    add_text_box(
        slide, MARGIN + Cm(1), Cm(1.0), CONTENT_W - Cm(2), Cm(5.5),
        "Genome-wide identification of candidate Ubiquitous Chromatin\n"
        "Opening Elements reveals a compositional nucleosome exclusion\n"
        "mechanism and conserved ETS binding motifs",
        font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(
        slide, MARGIN, Cm(6.8), CONTENT_W, Cm(2),
        "Elton Roger da Silva Ostetti\u00b9\u00b2, "
        "T\u00e2nia Maria Manieri\u00b9, "
        "Ana Maria Moro\u00b9",
        font_size=24, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB),
        alignment=PP_ALIGN.CENTER)

    add_text_box(
        slide, MARGIN, Cm(9.0), CONTENT_W, Cm(2.5),
        "\u00b9Laborat\u00f3rio de Biof\u00e1rmacos, Instituto Butantan, "
        "S\u00e3o Paulo, Brazil\n"
        "\u00b2Programa de P\u00f3s-gradua\u00e7\u00e3o Interunidades em "
        "Biotecnologia, Universidade de S\u00e3o Paulo, Brazil",
        font_size=18, color=RGBColor(0x90, 0xCA, 0xF9),
        alignment=PP_ALIGN.CENTER)

    add_text_box(
        slide, MARGIN, Cm(11.5), CONTENT_W, Cm(1.2),
        "Funded by FAPESP \u2022 S\u00e3o Paulo Research Foundation",
        font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)

    # ================================================================
    #  CONTENT
    # ================================================================
    content_top = Cm(14.8)
    col1 = MARGIN
    col2 = MARGIN + COL_W + COL_GAP
    col3 = MARGIN + 2 * (COL_W + COL_GAP)

    # ====================  COLUMN 1  ====================
    y = content_top

    # INTRODUCTION
    y = add_section_header(slide, col1, y, COL_W, "Introduction")
    h = add_body_text(slide, col1, y, COL_W,
        "Ubiquitous Chromatin Opening Elements (UCOEs) are regulatory DNA "
        "sequences that maintain constitutively open chromatin and prevent "
        "transgene silencing regardless of chromosomal integration site. "
        "In gene therapy, UCOEs protect therapeutic transgenes from "
        "epigenetic silencing; in biomanufacturing, they stabilize "
        "recombinant protein expression in producer cell lines.\n\n"
        "Only three UCOEs have been identified in the human genome: "
        "A2UCOE (HNRNPA2B1/CBX3), TBP/PSMB1, and SRF-UCOE "
        "(SURF1/SURF2). All share a common genomic architecture: CpG "
        "islands at bidirectional promoters between housekeeping genes, "
        "with constitutive histone marks and DNA hypomethylation.\n\n"
        "We developed a computational pipeline that exploits this shared "
        "architecture to identify novel UCOE candidates genome-wide, "
        "using multi-tissue epigenomic data to enforce the ubiquity "
        "criterion central to the UCOE definition.",
        font_size=16)
    y += h + Cm(0.4)

    # OBJECTIVES
    y = add_section_header(slide, col1, y, COL_W, "Objectives")
    h = add_body_text(slide, col1, y, COL_W,
        "1. Develop a two-phase pipeline for de novo UCOE candidate "
        "identification in the human genome (GRCh38), integrating "
        "epigenomic data from 11 ENCODE cell lines (~41 GB).\n\n"
        "2. Characterize intrinsic DNA sequence properties that "
        "distinguish UCOE regions from generic CpG islands.\n\n"
        "3. Propose a mechanistic model for UCOE-mediated chromatin "
        "opening based on epigenomic, biophysical, and evolutionary "
        "evidence.",
        font_size=16)
    y += h + Cm(0.4)

    # DATA
    y = add_section_header(slide, col1, y, COL_W, "Data")
    h = add_body_text(slide, col1, y, COL_W,
        "All data from ENCODE (11 cell lines representing diverse "
        "tissues) and GENCODE v44 annotations:",
        font_size=15)
    y += h + Cm(0.1)
    h = add_data_table(slide, col1, y, COL_W)
    y += h + Cm(0.1)
    h = add_body_text(slide, col1, y, COL_W,
        "Cell lines: GM12878, K562, HepG2, H1-hESC, HUVEC, HSMM, "
        "NHLF, NHEK, HMEC, IMR-90, A549",
        font_size=13, color=MED_GRAY)
    y += h + Cm(0.4)

    # METHODS
    y = add_section_header(slide, col1, y, COL_W, "Methods")

    y = add_subsection_header(slide, col1, y, COL_W,
                              "Phase I \u2014 Biological Filtering")
    h = add_body_text(slide, col1, y, COL_W,
        "Five sequential filters based on known UCOE properties:\n"
        "  F1. Bidirectional promoters between housekeeping genes\n"
        "  F2. CpG island overlap (\u226540%)\n"
        "  F3. Ubiquitous active marks (H3K4me3, H3K27ac; \u226580%)\n"
        "  F4. Absence of repressive marks (H3K27me3, H3K9me3)\n"
        "  F5. Constitutive hypomethylation (<10%; \u226580% of lines)",
        font_size=15)
    y += h + Cm(0.2)

    y = add_subsection_header(slide, col1, y, COL_W,
                              "Phase II \u2014 Multivariate Ranking")
    h = add_body_text(slide, col1, y, COL_W,
        "18 epigenomic features per candidate. Ranking by Mahalanobis "
        "distance and cosine similarity to three reference UCOEs. "
        "Sensitivity analysis: 29 weight combinations (Mahalanobis, "
        "Cosine, Percentile; 0.1\u20130.6 each, sum = 1.0).",
        font_size=15)
    y += h + Cm(0.2)

    y = add_subsection_header(slide, col1, y, COL_W,
                              "Sequence-level Analyses")
    h = add_body_text(slide, col1, y, COL_W,
        "All vs. 200 random CpG island controls (GC-matched). "
        "Mann-Whitney U + Benjamini-Hochberg FDR.\n"
        "  \u2022 DNA biophysics: nucleosome occupancy, stiffness\n"
        "  \u2022 Dinucleotide spectrum: WW/SS composition, periodicity\n"
        "  \u2022 K-mer enrichment (k=4,5,6; Fisher exact test)\n"
        "  \u2022 Conservation: PhyloP & PhastCons (100 vertebrates)",
        font_size=15)
    y += h + Cm(0.4)

    # Filter funnel
    fig_h = add_figure(slide, col1, y,
                       FIGURES_DIR / "fig2_filter_funnel.png",
                       width=COL_W, max_height=MAX_FIG_H)
    y += fig_h + Cm(0.1)
    add_caption(slide, col1, y, COL_W,
                "Figure 1. Phase I filtering funnel: 789 \u2192 599 "
                "candidates. All three known UCOEs recovered.")

    # ====================  COLUMN 2  ====================
    y = content_top

    # RESULTS
    y = add_section_header(slide, col2, y, COL_W, "Results")

    # Pipeline output
    y = add_subsection_header(slide, col2, y, COL_W, "Pipeline Output")
    h = add_body_text(slide, col2, y, COL_W,
        "Phase I recovered all three known human UCOEs and identified "
        "599 candidate regions across all autosomes and chrX. "
        "Phase II ranking followed by sensitivity analysis identified "
        "13 candidates with >80% stability, of which 7 reached 100%:",
        font_size=16)
    y += h + Cm(0.1)

    h = add_candidates_table(slide, col2, y, COL_W)
    y += h + Cm(0.1)
    h = add_caption(slide, col2, y, COL_W,
                    "Table 2. Candidates with 100% ranking stability "
                    "across all 29 weight combinations.")
    y += h + Cm(0.3)

    # Biophysics
    y = add_subsection_header(slide, col2, y, COL_W,
                              "Intrinsic DNA Biophysical Properties")
    h = add_body_text(slide, col2, y, COL_W,
        "UCOE candidates show significantly lower nucleosome formation "
        "potential (q = 5.76\u00d710\u207b\u00b9\u2076) and lower DNA "
        "stiffness (q = 7.38\u00d710\u207b\u2077) than CpG island "
        "controls. Nucleosome-enriched fraction drops from 84.7% to "
        "63.4% in candidates (\u221225%).",
        font_size=16)
    y += h + Cm(0.2)

    fig_h = add_figure(slide, col2, y,
                       FIGURES_DIR / "fig7_structural.png",
                       width=COL_W, max_height=MAX_FIG_H_SMALL)
    y += fig_h + Cm(0.1)
    h = add_caption(slide, col2, y, COL_W,
                    "Figure 2. Structural properties of UCOE candidates, "
                    "known UCOEs, and CpG island controls. Red stars: "
                    "known UCOEs.")
    y += h + Cm(0.3)

    # Dinucleotide signature
    y = add_subsection_header(slide, col2, y, COL_W,
                              "Dinucleotide Signature: Compositional, Not Periodic")
    h = add_body_text(slide, col2, y, COL_W,
        "Candidates contain 38% more WW dinucleotides (AA, AT, TA, TT; "
        "q = 7.82\u00d710\u207b\u00b2\u2070) and 9% fewer SS (CC, CG, "
        "GC, GG; q = 3.02\u00d710\u207b\u00b9\u00b9). The ~10.5 bp "
        "nucleosomal periodicity does NOT differ (q = 0.31). The DNA "
        "retains the helical bending code but lacks the stacking energy "
        "for stable nucleosome wrapping.",
        font_size=16)
    y += h + Cm(0.2)

    fig_h = add_figure(slide, col2, y,
                       FIGURES_DIR / "fig4_boxplots.png",
                       width=COL_W, max_height=MAX_FIG_H_SMALL)
    y += fig_h + Cm(0.1)
    h = add_caption(slide, col2, y, COL_W,
                    "Figure 3. Dinucleotide composition and periodicity: "
                    "WW fraction, SS fraction, autocorrelation at 10 bp, "
                    "entropy, Markov deviation, \u03c1*(CG).")
    y += h + Cm(0.3)

    # K-mer enrichment
    y = add_subsection_header(slide, col2, y, COL_W,
                              "K-mer Enrichment: ETS Binding Motif")
    h = add_body_text(slide, col2, y, COL_W,
        "K-mer analysis (k = 4, 5, 6) shows the most enriched motifs "
        "in candidates are A/T-rich, consistent with WW enrichment. "
        "The ETS (E26 transformation-specific) motif CGGAAG (fold "
        "enrichment = 1.94; q = 1.33\u00d710\u207b\u00b9\u00b3) is the "
        "only 6-mer simultaneously enriched in candidates AND present "
        "in all three known UCOEs. ETS factors GABPA and ELF1 are "
        "ubiquitous regulators of housekeeping gene promoters.",
        font_size=16)
    y += h + Cm(0.2)

    fig_h = add_figure(slide, col2, y,
                       FIGURES_DIR / "fig7_volcano_kmer.png",
                       width=COL_W, max_height=MAX_FIG_H)
    y += fig_h + Cm(0.1)
    add_caption(slide, col2, y, COL_W,
                "Figure 4. Volcano plot of 6-mer enrichment. Red: enriched "
                "+ shared across all known UCOEs.")

    # ====================  COLUMN 3  ====================
    y = content_top

    # RESULTS (continued)
    y = add_section_header(slide, col3, y, COL_W, "Results (continued)")

    y = add_subsection_header(slide, col3, y, COL_W,
                              "Evolutionary Conservation (100 Vertebrates)")
    h = add_body_text(slide, col3, y, COL_W,
        "Global conservation (PhyloP, PhastCons) does not differ between "
        "candidates and controls (q = 0.192). UCOEs are not exceptionally "
        "conserved as a class. However, within candidates, ETS motif "
        "positions show PhyloP = 0.901 vs. 0.371 outside motifs "
        "(p = 1.54\u00d710\u207b\u00b9\u2074), placing these non-coding "
        "motifs under purifying selection comparable to protein-coding "
        "exons.",
        font_size=16)
    y += h + Cm(0.2)

    fig_h = add_figure(slide, col3, y,
                       FIGURES_DIR / "fig11_motif_conservation.png",
                       width=COL_W, max_height=MAX_FIG_H_SMALL)
    y += fig_h + Cm(0.1)
    h = add_caption(slide, col3, y, COL_W,
                    "Figure 5. Left: PhyloP at ETS motifs vs. outside "
                    "(p = 1.54\u00d710\u207b\u00b9\u2074). "
                    "Right: ETS motif count distribution.")
    y += h + Cm(0.4)

    # PROPOSED MODEL
    y = add_section_header(slide, col3, y, COL_W, "Proposed Dual Mechanism")
    h = add_body_text(slide, col3, y, COL_W,
        "Four independent lines of evidence support a dual mechanism "
        "for UCOE-mediated chromatin opening:",
        font_size=16)
    y += h + Cm(0.1)

    h = add_highlight_box(slide, col3, y, COL_W,
        "1. INTRINSIC THERMODYNAMIC BARRIER\n"
        "WW-rich DNA composition lowers stacking energy, creating a "
        "thermodynamic barrier to stable nucleosome formation. The "
        "10.5 bp helical bending code is preserved, but nucleosomes "
        "cannot stabilize on these sequences.",
        font_size=15, bg_color=LIGHT_BLUE, text_color=DARK_BLUE)
    y += h + Cm(0.2)

    h = add_highlight_box(slide, col3, y, COL_W,
        "2. ACTIVE TRANSCRIPTION FACTOR COMPETITION\n"
        "Constitutive ETS factor occupancy (GABPA, ELF1) competes "
        "with transient nucleosomes. ETS motifs are enriched in all "
        "known UCOEs and in candidates, and are under strong purifying "
        "selection across 100 vertebrate genomes.",
        font_size=15, bg_color=RGBColor(0xE8, 0xF5, 0xE9),
        text_color=RGBColor(0x1B, 0x5E, 0x20))
    y += h + Cm(0.2)

    h = add_highlight_box(slide, col3, y, COL_W,
        "3. EPIGENETIC POSITIVE FEEDBACK\n"
        "Unmethylated CpGs recruit CXXC1/CFP1, depositing H3K4me3, "
        "which is mutually exclusive with DNA methylation \u2014 creating "
        "a self-reinforcing loop that sustains open chromatin.",
        font_size=15, bg_color=RGBColor(0xFD, 0xF1, 0xE1),
        text_color=RGBColor(0x7A, 0x4A, 0x00))
    y += h + Cm(0.2)

    h = add_highlight_box(slide, col3, y, COL_W,
        "4. EVOLUTIONARY VALIDATION\n"
        "ETS motifs within UCOE candidates are under purifying selection "
        "(PhyloP 2.4\u00d7 higher than flanking sequence), consistent "
        "with functional constraint maintained over vertebrate evolution.",
        font_size=15, bg_color=RGBColor(0xF3, 0xE5, 0xF5),
        text_color=RGBColor(0x4A, 0x14, 0x8C))
    y += h + Cm(0.5)

    # CONCLUSIONS
    y = add_section_header(slide, col3, y, COL_W, "Conclusions")
    h = add_body_text(slide, col3, y, COL_W,
        "\u2022 599 UCOE candidates identified from ~41 GB of ENCODE "
        "data (11 cell lines); 7 candidates stable under all 29 ranking "
        "weight configurations.\n\n"
        "\u2022 UCOE candidates exhibit a distinct DNA biophysical "
        "signature: lower nucleosome formation potential, lower stiffness, "
        "and a WW-rich dinucleotide composition that creates a "
        "thermodynamic barrier to nucleosome stabilization.\n\n"
        "\u2022 The ETS binding motif (CGGAAG) is enriched in candidates, "
        "shared across all known UCOEs, and under strong purifying "
        "selection \u2014 three independent lines of evidence "
        "supporting its functional role.\n\n"
        "\u2022 UCOEs maintain open chromatin through both intrinsic DNA "
        "topological properties (thermodynamic barrier) and active "
        "epigenetic mechanisms (ETS factor recruitment, CpG "
        "hypomethylation).\n\n"
        "\u2022 Top candidates are prioritized for experimental "
        "validation in transgene silencing assays.",
        font_size=16)
    y += h + Cm(0.5)

    # REFERENCES
    y = add_section_header(slide, col3, y, COL_W, "Key References")
    add_body_text(slide, col3, y, COL_W,
        "Antoniou, Harland & Mustoe (2003) Genomics 82:269\n"
        "ENCODE Consortium (2020) Nature 583:699\n"
        "Kohn et al. (2020) Nature Medicine 26:270\n"
        "Pollard et al. (2010) Genome Res 20:110\n"
        "Rudina & Smolke (2019) bioRxiv 626713\n"
        "Segal et al. (2006) Nature 442:772\n"
        "Siepel et al. (2005) Genome Res 15:1034\n"
        "Sizer & White (2022) CSBJ 21:275\n"
        "Williams et al. (2005) BMC Biotechnol 5:17",
        font_size=13, color=MED_GRAY)

    # ================================================================
    #  FOOTER
    # ================================================================
    footer_h = Cm(2)
    footer_top = H - footer_h
    add_colored_box(slide, Cm(0), footer_top, W, footer_h, DARK_BLUE)
    add_colored_box(slide, Cm(0), footer_top, W, Cm(0.15), GOLD)
    add_text_box(
        slide, MARGIN, footer_top + Cm(0.5), CONTENT_W, Cm(1.2),
        "X-Meeting 2026 \u2022 Brazilian Bioinformatics Congress \u2022 "
        "Contact: elton.ostetti@usp.br",
        font_size=18, color=RGBColor(0x90, 0xCA, 0xF9),
        alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTPUT))
    print(f"Poster saved to: {OUTPUT}")


if __name__ == "__main__":
    main()
